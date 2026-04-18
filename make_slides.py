from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Minecraft / SignQuest palette (from frontend) ─────────────────────────────
BG_DARK    = RGBColor(0x1A, 0x1A, 0x08)   # page background
BG_MID     = RGBColor(0x2D, 0x2D, 0x0F)   # mid tone
BG_CARD    = RGBColor(0x2A, 0x2A, 0x14)   # card background
BG_CARD2   = RGBColor(0x2A, 0x3A, 0x0A)   # badge bg
BORDER     = RGBColor(0x4A, 0x4A, 0x20)   # card border
GREEN_LT   = RGBColor(0xA0, 0xD0, 0x40)   # bright green (badge text)
GREEN_BTN  = RGBColor(0x5A, 0x7A, 0x1A)   # green button bg
GREEN_BDR  = RGBColor(0x8A, 0xB8, 0x28)   # green button border
GOLD       = RGBColor(0xC8, 0xA0, 0x30)   # gold accent
GOLD_TEXT  = RGBColor(0xF0, 0xC0, 0x30)   # gold text
PURPLE_BG  = RGBColor(0x5A, 0x1A, 0x7A)   # purple button
PURPLE_BDR = RGBColor(0x9A, 0x4A, 0xC8)   # purple border
TEAL_BG    = RGBColor(0x0A, 0x5A, 0x5A)   # teal button
TEAL_BDR   = RGBColor(0x20, 0xA8, 0xA8)   # teal border
STONE_BG   = RGBColor(0x5A, 0x5A, 0x5A)   # stone button
WHITE_TEXT = RGBColor(0xF0, 0xF0, 0xE0)   # main text
MUTED      = RGBColor(0xA0, 0xA0, 0x80)   # subtitle
DIM        = RGBColor(0x80, 0x80, 0x60)   # dim text
SHADOW     = RGBColor(0x3A, 0x3A, 0x10)   # text shadow sim
RED_WARN   = RGBColor(0xC8, 0x40, 0x40)   # warning red

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── helpers ───────────────────────────────────────────────────────────────────

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, border_color=None, border_pt=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color and border_pt:
        s.line.color.rgb = border_color
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size, color=WHITE_TEXT, bold=False,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb

def mc_badge(slide, label, l, t, w=4.0):
    """Green bordered badge like the frontend [ American Sign Language ] badge."""
    rect(slide, l, t, w, 0.42, BG_CARD2, GREEN_BDR, 2)
    txt(slide, label, l, t + 0.02, w, 0.38, 9, GREEN_LT, align=PP_ALIGN.CENTER)

def mc_card(slide, l, t, w, h, color=BG_CARD):
    """Blocky card with inset border."""
    rect(slide, l, t, w, h, color, BORDER, 3)

def section_title(slide, emoji_label, color=GREEN_LT):
    txt(slide, emoji_label, 0.5, 0.18, 12.3, 0.75, 28, color, bold=True)
    rect(slide, 0.5, 0.95, 12.3, 0.04, color)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
bg(slide)

# pixel-style background texture blocks (decorative)
for col, row in [(0,0),(12.5,0),(0,6.8),(12.5,6.8),(6,0.2)]:
    rect(slide, col, row, 0.6, 0.6, BG_CARD, BORDER, 1)

mc_badge(slide, "[ American Sign Language ]", 4.0, 1.1, 5.33)

txt(slide, "SignQuest", 1.5, 1.75, 10.33, 1.6, 64, WHITE_TEXT, bold=True, align=PP_ALIGN.CENTER)
# shadow effect (offset duplicate)
txt(slide, "SignQuest", 1.56, 1.82, 10.33, 1.6, 64, SHADOW, bold=True, align=PP_ALIGN.CENTER)

txt(slide, "Learn ASL. Play. Grow.", 1.5, 3.3, 10.33, 0.6, 16, MUTED, align=PP_ALIGN.CENTER, italic=True)

rect(slide, 3.5, 4.05, 6.33, 0.04, BORDER)

txt(slide, "A gamified web app for learning American Sign Language\nthrough real-time ML feedback, games, and an AI tutor.",
    1.5, 4.2, 10.33, 1.0, 11, DIM, align=PP_ALIGN.CENTER)

txt(slide, "Hackathon 2026  ·  Solo Project", 1.5, 5.4, 10.33, 0.5, 9, DIM, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
bg(slide)
section_title(slide, "😔  The Problem", RED_WARN)

# big stat card
mc_card(slide, 0.5, 1.15, 5.8, 2.1, RGBColor(0x2A, 0x10, 0x10))
txt(slide, "70%", 0.7, 1.2, 2.8, 1.3, 52, RED_WARN, bold=True)
txt(slide, "of deaf children in Nigeria\nare out of school", 3.3, 1.45, 2.8, 1.2, 12, WHITE_TEXT)

# story card
mc_card(slide, 0.5, 3.45, 5.8, 2.8, BG_CARD)
txt(slide, "A Personal Story", 0.7, 3.55, 5.4, 0.45, 10, GOLD_TEXT, bold=True)
txt(slide,
    "My cousin, who has hearing difficulties, was never taught\n"
    "sign language early on — leading to social isolation and\n"
    "a language barrier that could have been prevented.",
    0.7, 4.05, 5.4, 1.8, 11, MUTED)

# right panel
mc_card(slide, 6.8, 1.15, 6.0, 5.1, BG_CARD)
txt(slide, "The Reality", 7.0, 1.3, 5.6, 0.5, 12, GREEN_LT, bold=True)
points = [
    "Sign language is marginalized worldwide",
    "Severe communication barriers for the deaf",
    "Lack of accessible education for children",
    "Early language gaps cause lifelong impact",
    "Most families don't know where to start",
]
for i, p in enumerate(points):
    txt(slide, f"▸  {p}", 7.0, 1.95 + i * 0.82, 5.6, 0.7, 11, WHITE_TEXT if i % 2 == 0 else MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
bg(slide)
section_title(slide, "💡  The Solution", GREEN_LT)

txt(slide,
    "SignQuest is a gamified web app that teaches ASL letters & numbers\n"
    "through real-time ML-powered feedback — built for children.",
    0.5, 1.1, 12.3, 0.9, 12, MUTED)

cards_data = [
    ("🎮", "Gamified\nLearning",    GREEN_BTN,  GREEN_BDR,  "Earn points, unlock\nlevels, stay motivated"),
    ("📷", "Real-Time\nML Feedback", STONE_BG,  BORDER,     "Webcam model corrects\nyour signing live"),
    ("🤖", "AI Tutor",              PURPLE_BG,  PURPLE_BDR, "Recommends videos,\nclasses & schools"),
    ("🔁", "Spaced\nRepetition",    TEAL_BG,    TEAL_BDR,   "Signs stick through\nsmart review cycles"),
]
for i, (emoji, title, bg_c, bdr_c, desc) in enumerate(cards_data):
    x = 0.5 + i * 3.2
    rect(slide, x, 2.2, 3.0, 4.6, bg_c, bdr_c, 3)
    txt(slide, emoji,  x + 0.9, 2.4,  1.2, 0.9, 30, WHITE_TEXT, align=PP_ALIGN.CENTER)
    txt(slide, title,  x + 0.1, 3.35, 2.8, 0.8, 12, WHITE_TEXT, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, desc,   x + 0.1, 4.2,  2.8, 1.4, 10, RGBColor(0xD0,0xFF,0xD0), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Target Audience
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
bg(slide)
section_title(slide, "🎯  Who Is This For?", GOLD_TEXT)

# left card
mc_card(slide, 0.5, 1.15, 5.8, 5.8, BG_CARD)
txt(slide, "🧒", 2.3, 1.4, 2.2, 1.5, 56, WHITE_TEXT, align=PP_ALIGN.CENTER)
txt(slide, "Children with\nhearing disabilities", 0.7, 2.95, 5.2, 1.0, 14, GOLD_TEXT, bold=True, align=PP_ALIGN.CENTER)
txt(slide,
    "Early development is the critical\nwindow for language acquisition.\nWe meet them there.",
    0.7, 4.05, 5.2, 1.6, 11, MUTED, align=PP_ALIGN.CENTER)

# right card
mc_card(slide, 6.8, 1.15, 6.0, 5.8, BG_CARD)
txt(slide, "Why Start Early?", 7.0, 1.3, 5.6, 0.5, 12, GOLD_TEXT, bold=True)
rect(slide, 7.0, 1.85, 5.6, 0.04, BORDER)
why = [
    "Language skills form in early childhood",
    "Early ASL exposure improves literacy",
    "Prevents social isolation before it starts",
    "Games = higher engagement for kids",
    "Builds confidence and communication",
]
for i, w in enumerate(why):
    mc_card(slide, 7.0, 2.05 + i * 0.95, 5.6, 0.82, BG_MID)
    txt(slide, f"✓  {w}", 7.15, 2.12 + i * 0.95, 5.3, 0.68, 11, GREEN_LT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Competition
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
bg(slide)
section_title(slide, "⚔️  vs. The Competition", PURPLE_BDR)

# column headers
mc_card(slide, 0.5, 1.15, 5.8, 0.6, RGBColor(0x3A, 0x10, 0x3A))
mc_card(slide, 6.8, 1.15, 6.0, 0.6, GREEN_BTN)
txt(slide, "Others", 0.5, 1.2, 5.8, 0.5, 13, MUTED, bold=True, align=PP_ALIGN.CENTER)
txt(slide, "SignQuest  🤟", 6.8, 1.2, 6.0, 0.5, 13, WHITE_TEXT, bold=True, align=PP_ALIGN.CENTER)

rows = [
    ("Static flashcards only",          "Real-time webcam signing practice"),
    ("No feedback on accuracy",          "ML model corrects your signs live"),
    ("Passive watching",                 "Active gamified challenges"),
    ("No personalization",               "AI tutor + spaced repetition"),
    ("Not built for children",           "Designed for early childhood learning"),
]
for i, (them, us) in enumerate(rows):
    y = 1.9 + i * 0.95
    row_bg = BG_CARD if i % 2 == 0 else BG_MID
    mc_card(slide, 0.5, y, 5.8, 0.82, row_bg)
    mc_card(slide, 6.8, y, 6.0, 0.82, row_bg)
    txt(slide, f"✗  {them}", 0.65, y + 0.12, 5.5, 0.6, 11, RED_WARN)
    txt(slide, f"✓  {us}",   6.95, y + 0.12, 5.7, 0.6, 11, GREEN_LT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — How It Works
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
bg(slide)
section_title(slide, "⚙️  How It Works", TEAL_BDR)

steps = [
    ("📷", "Webcam\nCapture",    "MediaPipe\nhand landmarks"),
    ("🧠", "ML Classifier",     "Scikit-learn\nletters & numbers"),
    ("🎮", "Game Engine",       "React frontend\ngamified UI"),
    ("🤖", "AI Tutor Agent",    "LLM + spaced\nrepetition"),
]
for i, (emoji, title, sub) in enumerate(steps):
    x = 0.5 + i * 3.2
    mc_card(slide, x, 1.2, 2.9, 3.8, BG_CARD)
    txt(slide, emoji, x + 0.75, 1.4,  1.4, 1.0, 32, WHITE_TEXT, align=PP_ALIGN.CENTER)
    txt(slide, title, x + 0.1,  2.45, 2.7, 0.7, 12, GREEN_LT, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, sub,   x + 0.1,  3.2,  2.7, 1.2, 10, MUTED, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(slide, "→", x + 3.0, 2.5, 0.4, 0.6, 20, GOLD_TEXT, bold=True, align=PP_ALIGN.CENTER)

# stack badge
mc_badge(slide, "React + TypeScript  ·  FastAPI  ·  MediaPipe  ·  Scikit-learn  ·  Supabase  ·  WebSockets",
         0.5, 5.3, 12.33)

txt(slide, "Architecture", 0.5, 5.95, 12.3, 0.4, 9, DIM)
txt(slide,
    "Browser  →  WebSocket  →  FastAPI backend  →  MediaPipe landmark extraction  →  ML classifier  →  Response",
    0.5, 6.35, 12.3, 0.6, 9, DIM)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Team
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
bg(slide)
section_title(slide, "👤  The Team", GREEN_LT)

# left — solo dev card
mc_card(slide, 0.5, 1.2, 5.8, 5.5, BG_CARD)
txt(slide, "🧑‍💻", 2.3, 1.45, 2.2, 1.5, 56, WHITE_TEXT, align=PP_ALIGN.CENTER)
txt(slide, "Solo Founder", 0.7, 2.95, 5.2, 0.6, 16, GREEN_LT, bold=True, align=PP_ALIGN.CENTER)
rect(slide, 1.0, 3.62, 4.8, 0.04, BORDER)
txt(slide, "Full-stack dev  ·  ML engineer  ·  Designer", 0.7, 3.75, 5.2, 0.5, 10, MUTED, align=PP_ALIGN.CENTER)
txt(slide, "(yes, all of it 😅)", 0.7, 4.28, 5.2, 0.4, 10, DIM, align=PP_ALIGN.CENTER, italic=True)

# skills row
skills = ["React", "FastAPI", "MediaPipe", "Sklearn", "Supabase"]
for i, sk in enumerate(skills):
    x = 0.7 + i * 1.1
    mc_card(slide, x, 4.82, 1.0, 0.75, BG_MID)
    txt(slide, sk, x + 0.02, 4.9, 0.96, 0.58, 7, GREEN_LT, align=PP_ALIGN.CENTER)

# right — AI tools card
mc_card(slide, 7.0, 1.2, 5.8, 5.5, BG_CARD)
txt(slide, "🤖  AI Tools Used", 7.2, 1.35, 5.4, 0.55, 14, GOLD_TEXT, bold=True)
rect(slide, 7.2, 1.98, 5.4, 0.04, BORDER)
txt(slide, "Built with the help of AI assistants throughout\nthe entire development process.",
    7.2, 2.1, 5.4, 0.85, 10, MUTED)

ai_tools = [
    ("✦", "Kiro",   TEAL_BDR,   "AI IDE — code generation,\nspec writing & task planning"),
    ("✦", "Gemini", GOLD_TEXT,  "Research, architecture\ndesign & problem solving"),
    ("✦", "Claude", PURPLE_BDR, "Code review, debugging\n& refactoring assistance"),
]
for i, (dot, name, color, desc) in enumerate(ai_tools):
    y = 3.1 + i * 1.15
    mc_card(slide, 7.2, y, 5.4, 1.0, BG_MID)
    txt(slide, dot, 7.35, y + 0.18, 0.4, 0.65, 14, color, bold=True)
    txt(slide, name, 7.75, y + 0.1, 1.2, 0.45, 12, color, bold=True)
    txt(slide, desc, 7.75, y + 0.52, 4.6, 0.42, 9, MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
bg(slide)
section_title(slide, "🚀  What's Next", PURPLE_BDR)

phases = [
    ("Phase 1\n(Now ✓)",  GREEN_BTN,  GREEN_BDR,  [
        "ASL letters A–Z",
        "Numbers 0–9",
        "Gamified practice",
        "AI tutor agent",
    ]),
    ("Phase 2",           RGBColor(0x7A,0x60,0x10), GOLD, [
        "Full ASL word library",
        "LSTM model for words",
        "Translation feature",
        "More sign languages",
    ]),
    ("Phase 3",           TEAL_BG,    TEAL_BDR,   [
        "Full sentence recognition",
        "NLP context understanding",
        "Classroom tools",
        "Mobile app",
    ]),
]
for i, (phase, bg_c, bdr_c, items) in enumerate(phases):
    x = 0.5 + i * 4.3
    rect(slide, x, 1.15, 4.0, 5.7, bg_c, bdr_c, 3)
    txt(slide, phase, x + 0.1, 1.25, 3.8, 0.8, 13, WHITE_TEXT, bold=True, align=PP_ALIGN.CENTER)
    rect(slide, x, 2.1, 4.0, 0.04, bdr_c)
    for j, item in enumerate(items):
        txt(slide, f"▸  {item}", x + 0.2, 2.3 + j * 1.05, 3.6, 0.9, 11, WHITE_TEXT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — User Feedback
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
bg(slide)
section_title(slide, "💬  Real Feedback", GOLD_TEXT)

# centre quote card
mc_card(slide, 1.5, 1.2, 10.33, 3.8, BG_CARD)

# big quote mark
txt(slide, "\u201c", 1.8, 1.1, 1.2, 1.4, 72, GOLD_TEXT, bold=True)

txt(slide,
    "I actually understood what my hands were supposed to look like.\n"
    "The game made me want to keep going — I didn't want to stop!",
    2.2, 2.0, 9.0, 1.6, 14, WHITE_TEXT, italic=True)

txt(slide, "\u201d", 10.8, 2.8, 1.0, 1.0, 72, GOLD_TEXT, bold=True)

# attribution
rect(slide, 1.5, 4.85, 10.33, 0.04, BORDER)
txt(slide, "— My cousin, first-time user with hearing difficulties",
    1.5, 5.0, 10.33, 0.55, 11, GOLD_TEXT, italic=True, align=PP_ALIGN.CENTER)

# small context note
mc_card(slide, 3.5, 5.7, 6.33, 0.9, BG_MID)
txt(slide, "✓  Tested with a real user  ·  Positive first session  ·  Kept practicing independently",
    3.65, 5.82, 6.0, 0.65, 9, GREEN_LT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Call to Action
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
bg(slide)

# decorative corner blocks
for pos in [(0,0,1.5,1.5),(11.83,0,1.5,1.5),(0,6.0,1.5,1.5),(11.83,6.0,1.5,1.5)]:
    rect(slide, pos[0], pos[1], pos[2], pos[3], BG_CARD, BORDER, 2)

mc_badge(slide, "[ SignQuest ]", 4.67, 0.6, 4.0)

txt(slide, "🤟", 5.7, 1.3, 1.93, 1.5, 64, WHITE_TEXT, align=PP_ALIGN.CENTER)

txt(slide, "Every child deserves a voice.",
    1.0, 2.9, 11.33, 0.9, 28, GREEN_LT, bold=True, align=PP_ALIGN.CENTER)

txt(slide,
    "SignQuest makes sign language fun, accessible, and effective\nfor the children who need it most.",
    1.5, 3.9, 10.33, 1.0, 14, MUTED, align=PP_ALIGN.CENTER)

rect(slide, 4.0, 5.05, 5.33, 0.04, BORDER)

txt(slide, "Let's build a world where no child is left without language.",
    1.5, 5.2, 10.33, 0.7, 11, DIM, align=PP_ALIGN.CENTER, italic=True)

# ── save ──────────────────────────────────────────────────────────────────────
prs.save("SignQuest_v3.pptx")
print("SignQuest_v3.pptx saved!")
